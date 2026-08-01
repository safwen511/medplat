from __future__ import annotations

import json
from datetime import datetime, timezone
from pathlib import Path

import fitz
import pytest
from docling_core.types.doc import (
    BoundingBox,
    CoordOrigin,
    DocItemLabel,
    DoclingDocument,
    ProvenanceItem,
    Size,
    TableCell,
    TableData,
)
from docx import Document
from pptx import Presentation
from typer.testing import CliRunner

from ingestion.cli import app
from ingestion.hashing import sha256_file
from ingestion.normalization.models import (
    BlockType,
    DocumentType,
    LocationType,
    ProcessingReport,
    ProcessingStatus,
    TechnicalSuitability,
)
from ingestion.normalization.normalizer import normalize_docling_document
from ingestion.normalization.validation import validate_document_file, validation_summary
from ingestion.output import (
    OutputExistsError,
    UnsafeOutputError,
    write_normalized_output,
)
from ingestion.parsers.base import ParsedSource, StructuredDocumentParser, StructuredParserRegistry
from ingestion.parsers.docling_parser import (
    DoclingStructuredParser,
    UnsupportedFormatError,
    UnsupportedStructuredParser,
)
from ingestion.pipeline import parse_document, parse_sample


def make_pdf(path: Path) -> None:
    document = fitz.open()
    page = document.new_page(width=300, height=400)
    page.insert_text((30, 50), "Lesson title")
    page.insert_text((30, 120), "Clinical paragraph")
    document.save(path)
    document.close()


def make_low_text_image_pdf(path: Path) -> None:
    document = fitz.open()
    page = document.new_page(width=300, height=400)
    page.insert_text((30, 50), "tiny")
    pixmap = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 20, 20), False)
    pixmap.clear_with(255)
    page.insert_image(fitz.Rect(30, 80, 250, 300), pixmap=pixmap)
    document.save(path)
    document.close()


def provenance(page: int, top: float) -> ProvenanceItem:
    return ProvenanceItem(
        page_no=page,
        bbox=BoundingBox(
            l=30,
            t=top,
            r=200,
            b=top + 20,
            coord_origin=CoordOrigin.TOPLEFT,
        ),
        charspan=(0, 10),
    )


def structured_pdf_document() -> DoclingDocument:
    document = DoclingDocument(name="fixture")
    document.add_page(page_no=1, size=Size(width=300, height=400))
    title = document.add_title("Lesson title", prov=provenance(1, 30))
    heading = document.add_heading(
        "Clinical section", level=1, prov=provenance(1, 70), parent=title
    )
    document.add_text(
        DocItemLabel.PARAGRAPH,
        "Clinical paragraph",
        prov=provenance(1, 110),
        parent=heading,
    )
    data = TableData(
        num_rows=2,
        num_cols=2,
        table_cells=[
            TableCell(
                start_row_offset_idx=row,
                end_row_offset_idx=row + 1,
                start_col_offset_idx=column,
                end_col_offset_idx=column + 1,
                text=text,
                column_header=row == 0,
            )
            for row, values in enumerate((("Name", "Value"), ("HR", "72")))
            for column, text in enumerate(values)
        ],
    )
    document.add_table(data=data, prov=provenance(1, 160), parent=heading)
    subheading = document.add_heading("Details", level=2, prov=provenance(1, 220), parent=heading)
    document.add_text(
        DocItemLabel.PARAGRAPH,
        "Nested content",
        prov=provenance(1, 250),
        parent=subheading,
    )
    return document


def normalize_pdf(path: Path):  # type: ignore[no-untyped-def]
    return normalize_docling_document(
        structured_pdf_document(),
        source_path=path,
        source_relative_path="nested/lesson.pdf",
        sha256=sha256_file(path),
        mime_type="application/pdf",
        document_type=DocumentType.PDF,
        parser_name="docling-test",
        parser_version="1",
    )


def normalize_office(
    path: Path,
    document_type: DocumentType,
    *,
    text: str | None,
    initial_warnings: tuple[str, ...] = (),
):  # type: ignore[no-untyped-def]
    document = DoclingDocument(name="office-fixture")
    if document_type is DocumentType.POWERPOINT:
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(path)
        document.add_page(page_no=1, size=Size(width=300, height=400))
        if text is not None:
            document.add_title(text, prov=provenance(1, 30))
    else:
        source = Document()
        source.save(path)
        if text is not None:
            document.add_title(text)
    return normalize_docling_document(
        document,
        source_path=path,
        source_relative_path=path.name,
        sha256=sha256_file(path),
        mime_type=None,
        document_type=document_type,
        parser_name="docling-test",
        parser_version="1",
        initial_warnings=initial_warnings,
    )


def processing_report(status: ProcessingStatus = ProcessingStatus.PARTIAL_SUCCESS):
    now = datetime.now(timezone.utc)
    return ProcessingReport(
        parser_name="docling-test",
        parser_version="1",
        start_time=now,
        completion_time=now,
        status=status,
    )


def test_structured_registry_selects_supported_and_fallback() -> None:
    fallback = UnsupportedStructuredParser()
    parser = DoclingStructuredParser()
    registry = StructuredParserRegistry(fallback)
    registry.register(parser)

    assert registry.parser_for(".pdf") is parser
    assert registry.parser_for(".PPTX") is parser
    assert registry.parser_for(".doc") is fallback
    with pytest.raises(UnsupportedFormatError):
        fallback.parse(Path("legacy.doc"))


def test_pdf_normalization_page_sections_table_and_references(tmp_path: Path) -> None:
    source = tmp_path / "lesson.pdf"
    make_pdf(source)

    normalized = normalize_pdf(source)

    assert normalized.document_id == sha256_file(source)
    assert normalized.document_type is DocumentType.PDF
    assert normalized.page_or_slide_count == 1
    assert normalized.pages[0].number == 1
    assert normalized.pages[0].location_type is LocationType.PAGE
    assert normalized.pages[0].width == 300
    assert normalized.sections[1].parent_section_id == normalized.sections[0].section_id
    assert normalized.tables[0].row_count == 2
    assert normalized.tables[0].column_count == 2
    assert normalized.tables[0].cells[3].text == "72"
    paragraph = next(
        block for block in normalized.pages[0].blocks if block.text == "Clinical paragraph"
    )
    assert paragraph.block_type is BlockType.PARAGRAPH
    assert paragraph.page_or_slide_number == 1
    assert paragraph.source_reference.source_relative_path == "nested/lesson.pdf"
    assert paragraph.source_reference.block_id == paragraph.block_id
    assert paragraph.bounding_box is not None
    assert paragraph.bounding_box.coordinate_origin == "TOPLEFT"
    assert (
        normalized.metadata["extraction_quality"]["technical_suitability"] == "ready_with_warnings"
    )


def test_pdf_normalization_marks_low_text_image_pages_as_requiring_ocr(
    tmp_path: Path,
) -> None:
    source = tmp_path / "mixed.pdf"
    make_low_text_image_pdf(source)

    normalized = normalize_pdf(source)

    quality = normalized.metadata["extraction_quality"]
    assert quality["technical_suitability"] == "requires_ocr"
    assert quality["low_text_image_page_count"] == 1
    assert quality["low_text_image_pages"] == [1]
    assert quality["ocr_enabled"] is False
    assert any("ocr_required_but_disabled" in warning for warning in normalized.processing.warnings)


def test_pdf_normalization_uses_nonduplicating_text_fallback(tmp_path: Path) -> None:
    source = tmp_path / "backend-text.pdf"
    make_pdf(source)
    document = DoclingDocument(name="empty-layout")
    document.add_page(page_no=1, size=Size(width=300, height=400))

    normalized = normalize_docling_document(
        document,
        source_path=source,
        source_relative_path="nested/backend-text.pdf",
        sha256=sha256_file(source),
        mime_type="application/pdf",
        document_type=DocumentType.PDF,
        parser_name="docling-test",
        parser_version="1",
    )

    fallback = [
        block
        for block in normalized.pages[0].blocks
        if block.metadata.get("pymupdf_text_fallback") is True
    ]
    assert fallback
    assert all(block.block_type is BlockType.UNKNOWN for block in fallback)
    assert all(block.bounding_box is not None for block in fallback)
    assert all(
        block.source_reference.source_relative_path == "nested/backend-text.pdf"
        for block in fallback
    )
    assert normalized.metadata["extraction_quality"]["canonical_text_character_count"] > 0
    assert any("Docling supplied no text" in warning for warning in normalized.processing.warnings)


def test_pptx_docling_normalization_uses_slide_numbers(tmp_path: Path) -> None:
    path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Slide title"
    slide.placeholders[1].text = "Body text"
    presentation.save(path)

    parsed = DoclingStructuredParser().parse(path)
    normalized = normalize_docling_document(
        parsed.document,
        source_path=path,
        source_relative_path="slides.pptx",
        sha256=sha256_file(path),
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        document_type=DocumentType.POWERPOINT,
        parser_name=parsed.parser_name,
        parser_version=parsed.parser_version,
    )

    assert normalized.page_or_slide_count == 1
    assert normalized.pages[0].location_type is LocationType.SLIDE
    assert normalized.pages[0].number == 1
    assert all(block.page_or_slide_number == 1 for block in normalized.pages[0].blocks)


def test_nonempty_pptx_is_ready_for_chunking(tmp_path: Path) -> None:
    normalized = normalize_office(
        tmp_path / "nonempty.pptx", DocumentType.POWERPOINT, text="Slide title"
    )

    quality = normalized.metadata["extraction_quality"]
    assert quality["technical_suitability"] == "ready_for_chunking"
    assert quality["canonical_text_character_count"] == len("Slide title")


def test_warning_bearing_pptx_is_ready_with_warnings(tmp_path: Path) -> None:
    normalized = normalize_office(
        tmp_path / "warning.pptx",
        DocumentType.POWERPOINT,
        text="Slide title",
        initial_warnings=("Structured extraction requires review.",),
    )

    assert (
        normalized.metadata["extraction_quality"]["technical_suitability"] == "ready_with_warnings"
    )


def test_empty_pptx_is_unsuitable(tmp_path: Path) -> None:
    normalized = normalize_office(tmp_path / "empty.pptx", DocumentType.POWERPOINT, text=None)

    quality = normalized.metadata["extraction_quality"]
    assert quality["technical_suitability"] == "unsuitable"
    assert quality["canonical_text_character_count"] == 0


def test_docx_docling_normalization_has_null_pagination_and_table(tmp_path: Path) -> None:
    path = tmp_path / "notes.docx"
    source = Document()
    source.add_heading("Word title", 0)
    source.add_heading("Section", 1)
    source.add_paragraph("Paragraph")
    table = source.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    source.save(path)

    parsed = DoclingStructuredParser().parse(path)
    normalized = normalize_docling_document(
        parsed.document,
        source_path=path,
        source_relative_path="notes.docx",
        sha256=sha256_file(path),
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        document_type=DocumentType.WORD,
        parser_name=parsed.parser_name,
        parser_version=parsed.parser_version,
    )

    assert normalized.page_or_slide_count is None
    assert len(normalized.pages) == 1
    assert normalized.pages[0].location_type is LocationType.DOCUMENT
    assert normalized.pages[0].number is None
    assert all(block.page_or_slide_number is None for block in normalized.pages[0].blocks)
    assert normalized.tables[0].page_or_slide_number is None


def test_nonempty_docx_is_ready_with_pagination_warning(tmp_path: Path) -> None:
    normalized = normalize_office(
        tmp_path / "nonempty.docx", DocumentType.WORD, text="Document title"
    )

    assert normalized.page_or_slide_count is None
    assert normalized.pages[0].number is None
    assert (
        normalized.metadata["extraction_quality"]["technical_suitability"] == "ready_with_warnings"
    )


def test_warning_bearing_docx_is_ready_with_warnings(tmp_path: Path) -> None:
    normalized = normalize_office(
        tmp_path / "warning.docx",
        DocumentType.WORD,
        text="Document title",
        initial_warnings=("Structured extraction requires review.",),
    )

    assert (
        normalized.metadata["extraction_quality"]["technical_suitability"] == "ready_with_warnings"
    )
    assert "Structured extraction requires review." in normalized.processing.warnings


def test_empty_docx_is_unsuitable_and_keeps_null_pagination(tmp_path: Path) -> None:
    normalized = normalize_office(tmp_path / "empty.docx", DocumentType.WORD, text=None)

    quality = normalized.metadata["extraction_quality"]
    assert quality["technical_suitability"] == "unsuitable"
    assert quality["canonical_text_character_count"] == 0
    assert normalized.page_or_slide_count is None
    assert normalized.pages[0].number is None


def test_validation_and_summary(tmp_path: Path) -> None:
    source = tmp_path / "lesson.pdf"
    make_pdf(source)
    normalized = normalize_pdf(source)
    path = tmp_path / "document.json"
    path.write_text(normalized.model_dump_json(), encoding="utf-8")

    validated = validate_document_file(path)
    summary = validation_summary(validated)

    assert summary["schema_version"] == "1.0.0"
    assert summary["page_count"] == 1
    assert summary["slide_count"] is None
    assert summary["section_count"] == 2
    assert summary["table_count"] == 1


def test_atomic_write_and_overwrite_protection(tmp_path: Path) -> None:
    source_root = tmp_path / "library"
    source_root.mkdir()
    source = source_root / "lesson.pdf"
    make_pdf(source)
    normalized = normalize_pdf(source)
    output = tmp_path / "processed"

    final = write_normalized_output(
        normalized,
        processing_report(),
        source_path=source,
        source_root=source_root,
        output_root=output,
    )

    assert final == output / normalized.document_id
    assert (final / "document.json").is_file()
    assert (final / "document.md").is_file()
    assert (final / "processing-report.json").is_file()
    assert (final / "assets").is_dir()
    assert (final / "previews").is_dir()
    assert not list(output.glob("*.tmp"))
    with pytest.raises(OutputExistsError):
        write_normalized_output(
            normalized,
            processing_report(),
            source_path=source,
            source_root=source_root,
            output_root=output,
        )
    replacement = write_normalized_output(
        normalized,
        processing_report(),
        source_path=source,
        source_root=source_root,
        output_root=output,
        force=True,
    )
    assert replacement == final
    assert (
        validate_document_file(replacement / "document.json").document_id == normalized.document_id
    )


def test_atomic_failure_leaves_no_completed_or_temporary_result(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source_root = tmp_path / "library"
    source_root.mkdir()
    source = source_root / "lesson.pdf"
    make_pdf(source)
    normalized = normalize_pdf(source)
    output = tmp_path / "processed"

    def fail_validation(_path: Path):  # type: ignore[no-untyped-def]
        raise ValueError("validation failed")

    monkeypatch.setattr("ingestion.output.validate_document_file", fail_validation)
    with pytest.raises(ValueError, match="validation failed"):
        write_normalized_output(
            normalized,
            processing_report(),
            source_path=source,
            source_root=source_root,
            output_root=output,
        )

    assert not (output / normalized.document_id).exists()
    assert not list(output.glob(".*.tmp"))


def test_output_is_never_written_inside_source_directory(tmp_path: Path) -> None:
    source_root = tmp_path / "library"
    source_root.mkdir()
    source = source_root / "lesson.pdf"
    make_pdf(source)
    normalized = normalize_pdf(source)

    with pytest.raises(UnsafeOutputError):
        write_normalized_output(
            normalized,
            processing_report(),
            source_path=source,
            source_root=source_root,
            output_root=source_root / "generated",
        )
    assert not (source_root / "generated").exists()


class FixtureParser(StructuredDocumentParser):
    def __init__(self, fail_name: str | None = None) -> None:
        self.fail_name = fail_name

    @property
    def extensions(self) -> frozenset[str]:
        return frozenset({".docx"})

    def parse(self, path: Path) -> ParsedSource:
        if path.name == self.fail_name:
            raise RuntimeError("fixture failure")
        document = DoclingDocument(name="fixture")
        document.add_title("Fixture title")
        return ParsedSource(
            document,
            "fixture-parser",
            "1",
            provenance={"local_only": True, "ocr_enabled": False},
        )


def fixture_registry(parser: StructuredDocumentParser) -> StructuredParserRegistry:
    registry = StructuredParserRegistry(UnsupportedStructuredParser())
    registry.register(parser)
    return registry


def test_parse_document_is_content_addressed_and_preserves_relative_path(tmp_path: Path) -> None:
    source_root = tmp_path / "library"
    nested = source_root / "nested"
    nested.mkdir(parents=True)
    source = nested / "notes.docx"
    source.write_bytes(b"fixture")
    output = tmp_path / "processed"

    result = parse_document(
        source,
        source_root=source_root,
        output_root=output,
        registry=fixture_registry(FixtureParser()),
    )

    assert result.document.document_id == sha256_file(source)
    assert result.document.source_relative_path == "nested/notes.docx"
    assert result.output_directory.name == sha256_file(source)
    report = ProcessingReport.model_validate_json(
        (result.output_directory / "processing-report.json").read_text(encoding="utf-8")
    )
    assert report.model_provenance == {"local_only": True, "ocr_enabled": False}


def test_parse_document_downgrades_ready_suitability_when_warnings(tmp_path: Path) -> None:
    class WarningParser(FixtureParser):
        def parse(self, path: Path) -> ParsedSource:
            parsed = super().parse(path)
            return ParsedSource(
                parsed.document,
                parsed.parser_name,
                parsed.parser_version,
                warnings=("Structured extraction requires review.",),
                provenance=parsed.provenance,
            )

    source_root = tmp_path / "library"
    source_root.mkdir()
    source = source_root / "notes.docx"
    source.write_bytes(b"fixture")

    result = parse_document(
        source,
        source_root=source_root,
        output_root=tmp_path / "processed",
        registry=fixture_registry(WarningParser()),
        technical_suitability_override=TechnicalSuitability.READY_FOR_CHUNKING,
    )

    assert result.report.status is ProcessingStatus.PARTIAL_SUCCESS
    assert result.report.technical_suitability is TechnicalSuitability.READY_WITH_WARNINGS
    assert (
        result.document.metadata["extraction_quality"]["technical_suitability"]
        == "ready_with_warnings"
    )


def test_batch_failure_isolation_writes_report(tmp_path: Path) -> None:
    source_root = tmp_path / "library"
    source_root.mkdir()
    good = source_root / "good.docx"
    bad = source_root / "bad.docx"
    good.write_bytes(b"good")
    bad.write_bytes(b"bad")
    output = tmp_path / "processed"

    report = parse_sample(
        input_root=source_root,
        output_root=output,
        limit=2,
        registry=fixture_registry(FixtureParser(fail_name="bad.docx")),
        selected_paths=[bad, good],
    )

    assert [result.status for result in report.results] == [
        ProcessingStatus.FAILED,
        ProcessingStatus.PARTIAL_SUCCESS,
    ]
    payload = json.loads((output / "batch-processing-report.json").read_text(encoding="utf-8"))
    assert len(payload["results"]) == 2
    assert payload["results"][0]["status"] == "failed"


def test_validate_output_cli_exits_nonzero_for_invalid_json(tmp_path: Path) -> None:
    invalid = tmp_path / "document.json"
    invalid.write_text('{"schema_version": "wrong"}', encoding="utf-8")

    result = CliRunner().invoke(app, ["validate-output", str(invalid)])

    assert result.exit_code == 1
    assert "Validation failed" in result.stderr


def test_parse_sample_cli_requires_large_batch_acknowledgement(tmp_path: Path) -> None:
    result = CliRunner().invoke(
        app,
        ["parse-sample", "--input", str(tmp_path), "--limit", "4"],
    )

    assert result.exit_code == 2
    assert "--allow-large-batch" in result.stderr
