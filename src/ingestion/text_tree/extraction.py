"""Local format-specific extraction without OCR, network, or source writes."""

from __future__ import annotations

import codecs
import mimetypes
from dataclasses import dataclass
from pathlib import Path

from ingestion.config import DoclingSettings
from ingestion.models import DocumentClassification, FileMetadata
from ingestion.normalization.models import DocumentType
from ingestion.normalization.normalizer import normalize_docling_document
from ingestion.ocr.models import DocumentDerivative, OCRQualityOutcome
from ingestion.ocr.validation import validate_derivative_file
from ingestion.parsers.docling_parser import DoclingStructuredParser
from ingestion.parsers.pdf import PdfParser
from ingestion.text_tree.discovery import DiscoveredSource
from ingestion.text_tree.models import ExportStatus
from ingestion.text_tree.rendering import (
    RenderedBody,
    render_normalized_document,
    render_source_text,
)


@dataclass(frozen=True)
class PDFTechnicalInspection:
    """Minimum PDF evidence needed to plan extraction and OCR state."""

    classification: DocumentClassification
    readable: bool
    encrypted: bool | None
    character_count: int | None
    page_count: int | None
    error: str | None


@dataclass(frozen=True)
class AcceptedDerivative:
    """A validated accepted OCR derivative for an exact source identity."""

    metadata_path: Path
    pdf_path: Path
    derivative: DocumentDerivative


@dataclass(frozen=True)
class ExtractionResult:
    """One extraction outcome before metadata-header materialization."""

    status: ExportStatus
    extraction_tool: str | None
    document_type: str
    body: RenderedBody | None = None
    warnings: tuple[str, ...] = ()
    errors: tuple[str, ...] = ()


def inspect_pdf_technically(source: DiscoveredSource, input_root: Path) -> PDFTechnicalInspection:
    """Inspect PDF text availability without Docling initialization or OCR."""
    metadata = FileMetadata(
        relative_path=source.relative_path,
        filename=source.path.name,
        extension=source.extension,
        detected_type="pdf",
        mime_type="application/pdf",
        file_size=source.size_bytes,
        sha256=source.sha256,
    )
    inspected = PdfParser().inspect(source.path, metadata)
    error = None
    if inspected.error is not None:
        error = f"{inspected.error.error_type}: {inspected.error.message}"
    return PDFTechnicalInspection(
        classification=inspected.classification,
        readable=inspected.readable,
        encrypted=inspected.encrypted,
        character_count=inspected.extractable_character_count,
        page_count=inspected.page_count,
        error=error,
    )


def find_accepted_derivative(
    source: DiscoveredSource, *, input_root: Path, derived_root: Path
) -> tuple[AcceptedDerivative | None, tuple[str, ...]]:
    """Resolve only validated, explicitly accepted, exact-path/hash OCR derivatives."""
    if source.sha256 is None or not derived_root.is_dir():
        return None, ()
    warnings: list[str] = []
    matches: list[AcceptedDerivative] = []
    metadata_paths = sorted(
        derived_root.glob("*/ocr/*/derivative.json"), key=lambda path: path.as_posix()
    )
    accepted = {OCRQualityOutcome.ACCEPTED, OCRQualityOutcome.ACCEPTED_WITH_WARNINGS}
    for metadata_path in metadata_paths:
        try:
            candidate = DocumentDerivative.model_validate_json(
                metadata_path.read_text(encoding="utf-8")
            )
        except (OSError, ValueError):
            continue
        if (
            candidate.source_relative_path != source.relative_path
            or candidate.source_sha256 != source.sha256
            or candidate.validation_status not in accepted
        ):
            continue
        try:
            validated, _report = validate_derivative_file(metadata_path, source_root=input_root)
        except (OSError, ValueError) as exc:
            warnings.append(
                f"Matching accepted OCR derivative failed validation: {type(exc).__name__}: {exc}"
            )
            continue
        matches.append(
            AcceptedDerivative(
                metadata_path=metadata_path,
                pdf_path=metadata_path.parent / "document-ocr.pdf",
                derivative=validated,
            )
        )
    if not matches:
        return None, tuple(warnings)
    matches.sort(key=lambda item: item.derivative.derivative_id)
    if len(matches) > 1:
        warnings.append(
            "Multiple valid accepted OCR derivatives matched; selected the lowest derivative ID."
        )
    return matches[0], tuple(warnings)


def _pptx_notes(path: Path) -> tuple[dict[int, str], tuple[str, ...]]:
    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
        notes: dict[int, str] = {}
        for number, slide in enumerate(presentation.slides, start=1):
            if not slide.has_notes_slide:
                continue
            text = slide.notes_slide.notes_text_frame.text
            if text:
                notes[number] = text
        return notes, ()
    except Exception as exc:  # python-pptx exposes several package/XML errors.
        return {}, (f"Speaker notes were unavailable ({type(exc).__name__}).",)


def _decode_source_text(path: Path) -> tuple[str, str, tuple[str, ...]]:
    payload = path.read_bytes()
    if payload.startswith(codecs.BOM_UTF8):
        return payload.decode("utf-8-sig"), "utf-8-sig", ()
    if payload.startswith((codecs.BOM_UTF32_LE, codecs.BOM_UTF32_BE)):
        return payload.decode("utf-32"), "utf-32", ("Source encoding normalized to UTF-8.",)
    if payload.startswith((codecs.BOM_UTF16_LE, codecs.BOM_UTF16_BE)):
        return payload.decode("utf-16"), "utf-16", ("Source encoding normalized to UTF-8.",)
    return payload.decode("utf-8"), "utf-8", ()


class LocalTextExtractor:
    """Sequential reusable extractor with lazy local Docling initialization."""

    def __init__(self, docling_settings: DoclingSettings) -> None:
        self._parser = DoclingStructuredParser(settings=docling_settings)

    @property
    def pdf_initialized(self) -> bool:
        return self._parser.pdf_initialized

    def extract(
        self,
        source: DiscoveredSource,
        *,
        input_root: Path,
        pdf_inspection: PDFTechnicalInspection | None,
        accepted_derivative: AcceptedDerivative | None,
        derivative_warnings: tuple[str, ...] = (),
    ) -> ExtractionResult:
        """Extract one supported source while converting exceptions to failed states."""
        try:
            if source.extension == ".txt":
                value, encoding, text_warnings = _decode_source_text(source.path)
                body = render_source_text(value)
                if not value.strip():
                    return ExtractionResult(
                        status=ExportStatus.EMPTY,
                        extraction_tool=f"python-text-decoder/{encoding}",
                        document_type="text",
                    )
                status = (
                    ExportStatus.EXPORTED_WITH_WARNINGS if text_warnings else ExportStatus.EXPORTED
                )
                return ExtractionResult(
                    status=status,
                    extraction_tool=f"python-text-decoder/{encoding}",
                    document_type="text",
                    body=body,
                    warnings=text_warnings,
                )

            parse_path = source.path
            document_type = {
                ".pdf": DocumentType.PDF,
                ".pptx": DocumentType.POWERPOINT,
                ".docx": DocumentType.WORD,
            }[source.extension]
            structured_warnings = list(derivative_warnings)
            if source.extension == ".pdf":
                if pdf_inspection is None:
                    raise ValueError("PDF extraction requires a technical inspection.")
                if not pdf_inspection.readable or pdf_inspection.encrypted:
                    message = pdf_inspection.error or (
                        "Encrypted PDF cannot be extracted."
                        if pdf_inspection.encrypted
                        else "PDF is unreadable."
                    )
                    return ExtractionResult(
                        status=ExportStatus.FAILED,
                        extraction_tool="PyMuPDF inspection",
                        document_type="pdf",
                        errors=(message,),
                    )
                if (pdf_inspection.character_count or 0) == 0:
                    if accepted_derivative is None:
                        return ExtractionResult(
                            status=ExportStatus.REQUIRES_OCR,
                            extraction_tool="PyMuPDF inspection",
                            document_type="pdf",
                            warnings=tuple(structured_warnings),
                        )
                    parse_path = accepted_derivative.pdf_path
                    structured_warnings.append(
                        "Used a validated explicitly accepted OCR derivative with exact source "
                        "provenance."
                    )

            parsed = self._parser.parse(parse_path)
            structured_warnings.extend(parsed.warnings)
            mime_type, _encoding = mimetypes.guess_type(source.path.name)
            normalized = normalize_docling_document(
                parsed.document,
                source_path=parse_path,
                source_relative_path=source.relative_path,
                sha256=source.sha256 or "",
                mime_type=mime_type,
                document_type=document_type,
                parser_name=parsed.parser_name,
                parser_version=parsed.parser_version,
                initial_warnings=tuple(structured_warnings),
                source_filename=source.path.name,
            )
            notes: dict[int, str] = {}
            if source.extension == ".pptx":
                notes, note_warnings = _pptx_notes(source.path)
                structured_warnings.extend(note_warnings)
            structured_warnings.extend(
                warning
                for warning in normalized.processing.warnings
                if warning not in structured_warnings
            )
            body = render_normalized_document(normalized, notes_by_slide=notes)
            tool = str((parsed.provenance or {}).get("parser_backend", parsed.parser_name))
            if accepted_derivative is not None and parse_path == accepted_derivative.pdf_path:
                tool = f"{tool}+OCRmyPDF-accepted-derivative"
            if source.extension == ".pptx" and notes:
                tool = f"{tool}+python-pptx-notes"
            if body.text_character_count == 0:
                return ExtractionResult(
                    status=ExportStatus.EMPTY,
                    extraction_tool=tool,
                    document_type=document_type.value,
                    warnings=tuple(structured_warnings),
                )
            status = (
                ExportStatus.EXPORTED_WITH_WARNINGS
                if structured_warnings
                else ExportStatus.EXPORTED
            )
            return ExtractionResult(
                status=status,
                extraction_tool=tool,
                document_type=document_type.value,
                body=body,
                warnings=tuple(structured_warnings),
            )
        except (KeyboardInterrupt, SystemExit):
            raise
        except Exception as exc:
            return ExtractionResult(
                status=ExportStatus.FAILED,
                extraction_tool=None,
                document_type=source.extension.removeprefix(".") or "unknown",
                errors=(f"{type(exc).__name__}: {exc}",),
            )
